from smolagents import Tool
import easyocr
import os
import re
import json
import time
import torch
import requests
import pdfminer.high_level
import mammoth
import pandas as pd
import pptx
from smolagents.models import MessageRole, Model
import base64
from pathlib import Path
from typing import Optional, Dict, Any, List, Union
import subprocess
from translate import Translator
import puremagic  # 添加用于文件类型检测
import tempfile
import shutil

class FileProcessor:
    """文件处理的基础类，提供共享资源和通用功能"""
    
    def __init__(
        self,
        ocr_languages=["en", "ch_sim"],  # 使用 EasyOCR 支持的语言代码
        model=None  # 添加模型参数
    ):
        self.ocr_reader = None
        self.ocr_languages = ocr_languages
        self.model = model

    def detect_file_type(self, file_path: str) -> str:
        """检测文件类型并返回扩展名"""
        # 首先尝试从文件名获取扩展名
        ext = os.path.splitext(file_path)[1].lower()
        
        # 如果没有扩展名，使用 puremagic 进行检测
        if not ext:
            try:
                guesses = puremagic.magic_file(file_path)
                if guesses:
                    ext = "." + guesses[0].extension.strip()
            except (FileNotFoundError, IsADirectoryError, PermissionError):
                pass
                
        return ext

    def get_appropriate_tool(self, file_path: str) -> Optional[Tool]:
        """根据文件类型返回合适的处理工具"""
        ext = self.detect_file_type(file_path)
        
        # 文件类型到工具的映射
        tool_map = {
            '.pdf': PDFTool(self),
            '.docx': DOCXTool(self),
            '.xlsx': XLSXTool(self),
            '.xls': XLSXTool(self),
            '.pptx': PPTXTool(self)
        }
        
        return tool_map.get(ext)

    def get_ocr_reader(self):
        """初始化 OCR reader，使用正确的语言代码"""
        if self.ocr_reader is None:
            try:
                import easyocr
                self.ocr_reader = easyocr.Reader(self.ocr_languages, gpu=False)
            except Exception as e:
                raise Exception(f"OCR reader initialization failed: {str(e)}")
        return self.ocr_reader

class OCRTool(Tool):
    name = "OCR_Tool"
    description = "Extract text from images using OCR with automatic language detection."
    inputs = {
        "image_path": {
            "type": "string",
            "description": "Path to the image file.",
            "nullable": True
        },
        "languages": {
            "type": "array",
            "items": {
                "type": "string"
            },
            "description": "Optional list of language codes for OCR. If not provided, will auto-detect.",
            "default": ["en", "ch_sim"],
            "nullable": True
        },
        "auto_detect": {
            "type": "boolean",
            "description": "Whether to automatically detect languages in the image.",
            "default": True,
            "nullable": True
        }
    }
    output_type = "string"

    def __init__(self, file_processor: FileProcessor, model=None):
        super().__init__()
        self.file_processor = file_processor
        self.model = model  # LLM模型用于优化结果
        self._easyocr_readers = {}  # 缓存不同语言组合的读取器
        
        # 常见语言组合及其对应的代码
        self.language_groups = {
            "cjk": ["ch_sim", "ch_tra", "ja", "ko", "en"],  # 中日韩 + 英文
            "european": ["en", "fr", "de", "es", "it", "pt", "ru"],  # 欧洲主要语言
            "indic": ["hi", "ta", "te", "kn", "mr", "ne", "en"],  # 印度语系 + 英文
            "arabic": ["ar", "fa", "ur", "en"],  # 阿拉伯语系 + 英文
            "default": ["en", "ch_sim"]  # 默认组合
        }

    def _get_reader(self, languages=None):
        """获取或创建EasyOCR读取器"""
        if languages is None:
            languages = self.language_groups["default"]
            
        # 将语言列表转换为排序后的元组，用作字典键
        lang_key = tuple(sorted(languages))
        
        # 如果已经有缓存的读取器，则重用它
        if lang_key in self._easyocr_readers:
            return self._easyocr_readers[lang_key]
            
        # 否则创建新的读取器
        try:
            import easyocr
            print(f"初始化EasyOCR，支持语言: {languages}")
            # 修改：明确指定GPU=False，避免CUDA相关问题
            reader = easyocr.Reader(languages, gpu=False)
            self._easyocr_readers[lang_key] = reader
            return reader
        except Exception as e:
            print(f"初始化EasyOCR失败: {str(e)}")
            # 修改：返回默认读取器而不是None
            if lang_key != tuple(sorted(self.language_groups["default"])):
                print("尝试使用默认语言组合...")
                return self._get_reader(self.language_groups["default"])
            return None

    def _detect_language(self, image_path):
        """使用图像分析检测图像中的主要语言"""
        try:
            # 首先尝试使用基本的英文OCR来获取一些文本样本
            basic_reader = self._get_reader(["en"])
            if not basic_reader:
                return self.language_groups["default"]
                
            sample_text = basic_reader.readtext(image_path, detail=0)
            
            # 如果没有检测到文本，返回默认语言组合
            if not sample_text:
                return self.language_groups["default"]
                
            # 分析文本特征来猜测语言
            text = " ".join(sample_text)
            
            # 检测CJK字符（中日韩）
            if any(ord(c) > 0x4E00 and ord(c) < 0x9FFF for c in text):
                return self.language_groups["cjk"]
                
            # 检测西里尔字符（俄语等）
            if any(ord(c) > 0x0400 and ord(c) < 0x04FF for c in text):
                return ["ru", "en"]
                
            # 检测阿拉伯字符
            if any(ord(c) > 0x0600 and ord(c) < 0x06FF for c in text):
                return self.language_groups["arabic"]
                
            # 检测印度语系字符
            if any(ord(c) > 0x0900 and ord(c) < 0x097F for c in text):
                return self.language_groups["indic"]
                
            # 如果没有特殊字符，假设是欧洲语言
            return self.language_groups["european"]
            
        except Exception as e:
            print(f"语言检测失败: {str(e)}")
            return self.language_groups["default"]

    def forward(self, image_path: str = None, languages: List[str] = None, auto_detect: bool = True) -> str:
        try:
            if image_path is None:
                return "错误：未提供图像文件路径"
                
            if not os.path.exists(image_path):
                return f"错误：未找到文件路径: {image_path}"
                
            # 检查文件扩展名
            ext = os.path.splitext(image_path)[1].lower()
            if ext not in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.webp']:
                return "错误：不支持的图像格式。支持的格式包括：JPG, JPEG, PNG, BMP, TIFF, WEBP"
            
            # 如果未提供语言且启用了自动检测，则检测语言
            if languages is None and auto_detect:
                print("正在自动检测图像语言...")
                languages = self._detect_language(image_path)
                print(f"检测到可能的语言组合: {languages}")
            elif languages is None:
                languages = self.language_groups["default"]
                
            # 获取OCR读取器
            reader = self._get_reader(languages)
            if reader is None:
                # 修改：返回更详细的错误信息，并尝试使用系统描述
                print("无法初始化OCR引擎，尝试使用系统工具描述图像...")
                if self.model:
                    return self._describe_image_with_llm(image_path)
                return "错误：无法初始化OCR引擎，请检查EasyOCR安装和语言支持"
                
            # 执行OCR
            print(f"正在对图像 {image_path} 执行OCR，使用语言: {languages}")
            try:
                # 添加超时处理
                results = reader.readtext(image_path, detail=0)  # detail=0只返回文本
            except Exception as ocr_error:
                print(f"OCR处理失败: {str(ocr_error)}")
                # 尝试使用默认语言
                if languages != self.language_groups["default"]:
                    print("尝试使用默认语言组合...")
                    reader = self._get_reader(self.language_groups["default"])
                    if reader:
                        try:
                            results = reader.readtext(image_path, detail=0)
                        except:
                            results = []
                    else:
                        results = []
                else:
                    results = []
            
            # 检查结果
            if not results:
                print("OCR未检测到文本，尝试使用其他语言组合...")
                # 尝试使用其他语言组合
                for group_name, group_langs in self.language_groups.items():
                    if group_langs != languages:
                        print(f"尝试使用{group_name}语言组合: {group_langs}")
                        alt_reader = self._get_reader(group_langs)
                        if alt_reader:
                            try:
                                alt_results = alt_reader.readtext(image_path, detail=0)
                                if alt_results:
                                    results = alt_results
                                    print(f"使用{group_name}语言组合成功检测到文本")
                                    break
                            except Exception as e:
                                print(f"使用{group_name}语言组合OCR失败: {str(e)}")
                
                # 如果仍然没有结果，尝试使用LLM描述图像
                if not results and self.model:
                    return self._describe_image_with_llm(image_path)
                elif not results:
                    return "未能从图像中检测到任何文本，请尝试使用其他工具或手动处理"
                
            # 将结果组合成文本
            text = "\n".join(results)
            
            # 使用LLM优化OCR结果（如果有模型）
            if self.model and self._needs_optimization(text):
                return self._optimize_with_llm(text, image_path)
            
            return text
            
        except Exception as e:
            print(f"OCR处理失败: {str(e)}")
            # 如果有模型，尝试使用LLM描述图像
            if self.model:
                return self._describe_image_with_llm(image_path)
            return f"OCR处理失败: {str(e)}"

    def _needs_optimization(self, text: str) -> bool:
        """判断OCR结果是否需要优化"""
        if not text:
            return False
            
        # 检查文本是否包含特殊字符比例过高
        special_chars = sum(1 for c in text if not c.isalnum() and not c.isspace())
        if special_chars / len(text) > 0.3:  # 如果特殊字符超过30%
            return True
            
        # 检查是否有太多连续的非单词字符（可能是识别错误）
        if re.search(r'[^\w\s]{3,}', text):
            return True
            
        # 检查是否有太多不常见的字符组合
        if re.search(r'[a-z][A-Z]{2,}|[A-Z][a-z]{2,}[A-Z]', text):
            return True
            
        return True
        
    def _optimize_with_llm(self, text: str, image_path: str) -> str:
        """使用LLM优化OCR结果"""
        try:
            # 使用与TextInspectorTool相同的消息格式
            messages = [
                {
                    "role": MessageRole.SYSTEM,
                    "content": [
                        {
                            "type": "text",
                            "text": "你是一个OCR文本修正专家。你的任务是修复OCR识别错误，提供准确的文本。只返回修正后的文本，不要添加任何解释、分析或评论。"
                        }
                    ],
                },
                {
                    "role": MessageRole.USER,
                    "content": [
                        {
                            "type": "text",
                            "text": f"以下是OCR识别的原始文本，可能包含错误:\n\n{text}\n\n请直接提供修正后的纯文本，无需任何额外内容:"
                        }
                    ],
                }
            ]
            
            print("DEBUG: 使用与TextInspectorTool一致的调用方式")
            
            # 直接调用模型并获取content属性
            response = self.model(messages).content
            
            # 确保返回非空内容
            if response and response.strip():
                return response.strip()
            else:
                print("DEBUG: 响应为空，返回原始文本")
                return text
            
        except Exception as e:
            print(f"LLM优化失败: {str(e)}")
            # 确保处理所有可能的异常
            return text

    # 添加新方法：使用LLM描述图像
    def _describe_image_with_llm(self, image_path: str) -> str:
        """当OCR失败时，使用LLM描述图像内容"""
        try:
            # 读取图像文件并转换为base64
            with open(image_path, "rb") as image_file:
                image_data = base64.b64encode(image_file.read()).decode('utf-8')
            
            # 准备消息
            messages = [
                {
                    "role": MessageRole.SYSTEM,
                    "content": [
                        {
                            "type": "text",
                            "text": "你是一个图像分析专家。请详细描述图像中的内容，特别关注任何文本、符号、图表或重要视觉元素。"
                        }
                    ],
                },
                {
                    "role": MessageRole.USER,
                    "content": [
                        {
                            "type": "text",
                            "text": "请详细描述这张图像中的内容，特别是任何可见的文本:"
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{image_data}"
                            }
                        }
                    ],
                }
            ]
            
            # 调用模型
            response = self.model(messages).content
            
            if response and response.strip():
                return f"[图像描述] {response.strip()}"
            else:
                return "无法分析图像内容"
            
        except Exception as e:
            print(f"图像描述失败: {str(e)}")
            return f"无法处理图像: {str(e)}"

class FileProcessingResult:
    """文件处理结果的标准化格式"""
    def __init__(self, title: Optional[str] = None, content: str = "", error: Optional[str] = None):
        self.title = title
        self.content = content
        self.error = error

    def __str__(self):
        if self.error:
            return f"错误: {self.error}"
        if self.title:
            return f"{self.title}\n\n{self.content}"
        return self.content

class FileProcessingException(Exception):
    """文件处理异常的基类"""
    pass

class UnsupportedFormatException(FileProcessingException):
    """不支持的文件格式异常"""
    pass

class FileConversionException(FileProcessingException):
    """文件转换失败异常"""
    pass

class PDFTool(Tool):
    name = "pdf_tool"
    description = "Extract text from PDF files and perform analysis."
    inputs = {
        "file_path": {
            "type": "string",
            "description": "Path to the PDF file",
            "nullable": True
        },
        "page_range": {
            "type": "string", 
            "description": "Page range to extract (e.g., '1-5', 'all')",
            "default": "all",
            "nullable": True
        }
    }
    output_type = "string"

    def __init__(self, file_processor: FileProcessor):
        super().__init__()
        self.file_processor = file_processor
        
    def forward(self, file_path: str = None, page_range: str = "all") -> str:
        """Extract text from a PDF file."""
        if not file_path or not os.path.exists(file_path):
            return f"错误：文件不存在或路径无效 - {file_path}"
            
        try:
            # 尝试常规PDF文本提取
            print(f"尝试从PDF提取文本: {file_path}")
            text = ""
            
            # 确定页面范围
            pages = None
            if page_range != "all":
                try:
                    if "-" in page_range:
                        start, end = map(int, page_range.split("-"))
                        pages = list(range(start-1, end))  # pdfminer使用0索引
                    else:
                        pages = [int(page_range) - 1]
                except ValueError:
                    return f"错误：无效的页面范围 - {page_range}"
            
            # 第一次尝试：使用pdfminer
            try:
                print("使用pdfminer提取文本...")
                text = pdfminer.high_level.extract_text(file_path, page_numbers=pages)
            except Exception as e:
                print(f"pdfminer提取失败: {e}")
                return f"错误：PDF处理失败 - {str(e)}"
                
            # 检查是否成功提取文本
            if text and len(text.strip()) > 50:  # 至少有一些有意义的文本
                print(f"成功提取文本，长度: {len(text)} 字符")
                return format_pdf_text(text, file_path)
                
            # 第二次尝试：使用PyPDF2
            print("pdfminer提取失败或提取文本太少，尝试PyPDF2...")
            import PyPDF2
            text = ""
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                if pages is None:
                    pages = range(len(reader.pages))
                for i in pages:
                    if i < len(reader.pages):
                        page_text = reader.pages[i].extract_text()
                        if page_text:
                            text += f"\n--- 第 {i+1} 页 ---\n" + page_text
            
            # 检查第二次尝试结果
            if text and len(text.strip()) > 50:
                print(f"PyPDF2成功提取文本，长度: {len(text)} 字符")
                return format_pdf_text(text, file_path)
                
            # 如果两种方法都失败了，可能是扫描的PDF，尝试OCR
            print("直接文本提取失败，PDF可能是扫描件，尝试OCR...")
            
            # 检查OCR工具是否可用
            if hasattr(self.file_processor, 'ocr_reader'):
                # 使用PDF转图像然后OCR的方法
                from pdf2image import convert_from_path
                import tempfile
                
                print(f"将PDF转换为图像进行OCR处理...")
                with tempfile.TemporaryDirectory() as temp_dir:
                    # 确定要处理的页面
                    if pages is None:
                        pdf_pages = convert_from_path(file_path)
                    else:
                        pdf_pages = convert_from_path(file_path, first_page=min(pages)+1, last_page=max(pages)+1)
                    
                    all_text = []
                    for i, page in enumerate(pdf_pages):
                        # 保存图像
                        image_path = os.path.join(temp_dir, f'page_{i+1}.png')
                        page.save(image_path, 'PNG')
                        
                        # OCR处理
                        print(f"OCR处理第 {i+1} 页...")
                        result = self.file_processor.ocr_reader.readtext(image_path)
                        page_text = "\n".join([text for _, text, _ in result])
                        all_text.append(f"\n--- 第 {i+1} 页 (OCR) ---\n{page_text}")
                    
                    if all_text:
                        ocr_text = "\n".join(all_text)
                        print(f"OCR成功提取文本，长度: {len(ocr_text)} 字符")
                        return format_pdf_text(ocr_text, file_path, is_ocr=True)
            
            # 如果所有方法都失败
            return f"警告：无法从PDF提取文本。此PDF可能是扫描件、图像化文档或有特殊保护。请上传PDF截图以查看内容。"
            
        except Exception as e:
            import traceback
            trace = traceback.format_exc()
            return f"错误：PDF处理失败 - {str(e)}\n\n详细错误：\n{trace}"

def format_pdf_text(text, file_path, is_ocr=False):
    """格式化提取的PDF文本"""
    file_name = os.path.basename(file_path)
    method = "OCR" if is_ocr else "直接提取"
    
    # 获取文本统计信息
    text_length = len(text)
    lines = text.count('\n') + 1
    words = len(text.split())
    
    header = f"""# PDF文件分析：{file_name}
- 提取方法: {method}
- 文本长度: {text_length} 字符
- 行数: {lines}
- 词数: {words}

## 提取内容:
"""
    
    return header + text

class DOCXTool(Tool):
    name = "DOCX_Tool"
    description = "Convert DOCX files to HTML."
    inputs = {
        "file_path": {
            "type": "string",
            "description": "Path to the DOCX file.",
            "nullable": True
        }
    }
    output_type = "string"

    def __init__(self, file_processor: FileProcessor):
        super().__init__()
        self.file_processor = file_processor

    def forward(self, file_path: str = None) -> str:
        try:
            if file_path is None:
                return "错误：未提供DOCX文件路径"
                
            if not os.path.exists(file_path):
                return f"Error: File not found at path: {file_path}"
                
            # 检查文件扩展名
            if not file_path.lower().endswith('.docx'):
                return "Error: File is not a DOCX document."
                
            with open(file_path, "rb") as docx_file:
                result = mammoth.convert_to_html(docx_file)
                html_content = result.value
                
                if not html_content.strip():
                    return "No content found in the document."
                    
                return html_content.strip()
        except Exception as e:
            return f"Error processing DOCX: {str(e)}"

class XLSXTool(Tool):
    name = "XLSX_Tool"
    description = "Convert XLSX files to Markdown."
    inputs = {
        "file_path": {
            "type": "string",
            "description": "Path to the XLSX file.",
            "nullable": True
        }
    }
    output_type = "string"

    def __init__(self, file_processor: FileProcessor):
        super().__init__()
        self.file_processor = file_processor

    def forward(self, file_path: str = None) -> str:
        try:
            if file_path is None:
                return "错误：未提供Excel文件路径"
                
            if not os.path.exists(file_path):
                return f"Error: File not found at path: {file_path}"
                
            # 检查文件扩展名
            if not file_path.lower().endswith(('.xlsx', '.xls')):
                return "Error: File is not an Excel spreadsheet."
                
            sheets = pd.read_excel(file_path, sheet_name=None)
            if not sheets:
                return "No data found in the Excel file."
                
            md_content = ""
            for sheet_name, sheet_data in sheets.items():
                if sheet_data.empty:
                    continue
                    
                md_content += f"## {sheet_name}\n"
                md_content += sheet_data.to_markdown(index=False) + "\n\n"
                
            if not md_content:
                return "No data found in the Excel file."
                
            return md_content.strip()
        except Exception as e:
            return f"Error processing XLSX: {str(e)}"

class PPTXTool(Tool):
    name = "pptx_tool"
    description = "Extract text and structure from PowerPoint presentations, including image analysis."
    inputs = {
        "file_path": {
            "type": "string",
            "description": "Path to the PPTX file",
            "nullable": True
        }
    }
    output_type = "string"

    def __init__(self, file_processor: FileProcessor):
        super().__init__()
        self.file_processor = file_processor
        self.image_analyzer = None
        
    def forward(self, file_path: str = None) -> str:
        try:
            if not file_path or not os.path.exists(file_path):
                return f"错误：文件不存在或路径无效 - {file_path}"
                
            # 验证文件类型
            if self.file_processor.detect_file_type(file_path) != '.pptx':
                return "错误：文件不是PPTX格式"
                
            # 提取PPT内容
            presentation = pptx.Presentation(file_path)
            content = []
            
            # 添加标题
            content.append(f"# PowerPoint: {os.path.basename(file_path)}")
            
            # 创建临时目录用于保存图片
            temp_dir = tempfile.mkdtemp()
            
            try:
                # 初始化图像分析器（如果可能）
                self._init_image_analyzer()
                
                # 处理每个幻灯片
                slide_num = 0
                for slide in presentation.slides:
                    slide_num += 1
                    content.append(f"\n## 幻灯片 {slide_num}")
                    
                    # 处理幻灯片标题
                    title = slide.shapes.title
                    if title and title.has_text_frame:
                        content.append(f"### {title.text.strip()}")
                    
                    # 处理幻灯片中的各种形状
                    for shape in slide.shapes:
                        # 跳过已处理的标题
                        if shape == title:
                            continue
                            
                        # 处理图片
                        if self._is_picture(shape):
                            # 尝试获取alt文本
                            alt_text = ""
                            try:
                                alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
                            except Exception:
                                pass
                            
                            # 提取并保存图片
                            image_path = self._extract_image(shape, temp_dir, slide_num)
                            if image_path:
                                # 使用ImageAnalysisTool分析图片
                                image_description = self._analyze_image(image_path)
                                
                                # 添加图片描述
                                if alt_text:
                                    content.append(f"[图片描述: {alt_text}]")
                                content.append(f"[图片分析: {image_description}]")
                            else:
                                content.append("[图片: 无法提取]")
                        
                        # 处理表格
                        elif self._is_table(shape):
                            content.append("\n#### 表格内容:")
                            table_content = []
                            
                            # 表头
                            if shape.table.rows:
                                header = []
                                for cell in shape.table.rows[0].cells:
                                    header.append(cell.text.strip())
                                table_content.append("| " + " | ".join(header) + " |")
                                
                                # 分隔行
                                table_content.append("| " + " | ".join(["---"] * len(header)) + " |")
                                
                                # 表格内容
                                for row in shape.table.rows[1:]:
                                    row_content = []
                                    for cell in row.cells:
                                        row_content.append(cell.text.strip())
                                    table_content.append("| " + " | ".join(row_content) + " |")
                                    
                                content.append("\n".join(table_content))
                        
                        # 处理文本框
                        elif shape.has_text_frame:
                            text = shape.text.strip()
                            if text:
                                content.append(text)
                    
                    # 处理备注
                    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                        notes = slide.notes_slide.notes_text_frame.text.strip()
                        if notes:
                            content.append("\n#### 备注:")
                            content.append(notes)
            
            finally:
                # 清理临时目录
                shutil.rmtree(temp_dir, ignore_errors=True)
                
            return "\n\n".join(content)
            
        except Exception as e:
            return f"错误：PPTX处理失败 - {str(e)}"
    
    def _init_image_analyzer(self):
        """初始化图像分析器"""
        if self.image_analyzer is not None:
            return  # 已经初始化过了
            
        try:
            # 检查FileProcessor是否有模型
            if hasattr(self.file_processor, 'model') and self.file_processor.model is not None:
                # 创建图像分析器
                self.image_analyzer = ImageAnalysisTool(self.file_processor, self.file_processor.model)
            else:
                print("警告: FileProcessor没有可用的模型，无法进行图像分析")
        except Exception as e:
            print(f"初始化图像分析器失败: {str(e)}")
    
    def _is_picture(self, shape):
        """检查形状是否为图片"""
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
            try:
                if hasattr(shape, "image"):
                    return True
            except:
                pass
        return False
        
    def _is_table(self, shape):
        """检查形状是否为表格"""
        return hasattr(shape, "table")
        
    def _extract_image(self, shape, temp_dir, slide_num):
        """提取图片并保存到临时目录"""
        try:
            # 提取图片
            if hasattr(shape, "image"):
                # 从形状中获取图片
                image_bytes = shape.image.blob
                image_ext = shape.image.ext
                image_filename = f"slide_{slide_num}_{shape.name}.{image_ext}"
                image_path = os.path.join(temp_dir, image_filename)
                
                with open(image_path, "wb") as f:
                    f.write(image_bytes)
                    
                return image_path
            return None
        except Exception as e:
            print(f"图片提取失败: {str(e)}")
            return None
            
    def _analyze_image(self, image_path):
        """使用ImageAnalysisTool分析图片"""
        try:
            # 检查图像分析器是否可用
            if self.image_analyzer is None:
                return "无法分析图片：图像分析器未初始化"
                
            # 使用ImageAnalysisTool分析图片
            result = self.image_analyzer.forward(image_path)
            
            # 简化结果（移除标题等）
            if result and isinstance(result, str):
                # 移除可能的标题和前缀
                result = re.sub(r'^.*?\[图像描述\]\s*', '', result, flags=re.DOTALL)
                return result.strip()
            
            return "无法分析图片内容"
            
        except Exception as e:
            print(f"图像分析失败: {str(e)}")
            return f"图像分析失败: {str(e)}"

class ImageAnalysisTool(Tool):
    name = "Image_Analysis_Tool"
    description = "分析图片内容，提供详细的图像描述。输入图片路径，输出对图像的全面分析。"
    inputs = {
        "image_path": {
            "type": "string",
            "description": "Path to the image file.",
            "nullable": True
        }
    }
    output_type = "string"

    def __init__(self, file_processor: FileProcessor, model):
        """
        初始化图像分析工具
        
        Args:
            file_processor: 文件处理器实例
            model: 用于分析图像的模型
        """
        super().__init__()
        self.file_processor = file_processor
        self.model = model

    def forward(self, image_path: str = None) -> str:
        """
        对图像进行分析，返回模型生成的描述
        
        Args:
            image_path: 图像文件路径
        
        Returns:
            str: 图像内容的详细描述
        """
        try:
            if image_path is None:
                return "错误：未提供图像文件路径"
                
            # 检查文件是否存在
            if not os.path.exists(image_path):
                return f"错误：文件 {image_path} 不存在"
                
            # 获取文件扩展名
            ext = os.path.splitext(image_path)[1].lower()
            
            # 检查是否为支持的图像格式
            if ext not in ['.jpg', '.jpeg', '.png', '.bmp', '.gif', '.webp']:
                return f"错误：不支持的图像格式 {ext}"
                
            # 读取图像文件并进行base64编码
            with open(image_path, "rb") as image_file:
                encoded_image = base64.b64encode(image_file.read()).decode('utf-8')
                
            # 构建消息内容
            messages = [
                {
                    "role": "system",
                    "content": [
                        {
                            "type": "text",
                            "text": "你是一位专业的图像分析助手。请详细描述图像中的内容，包括可见的物体、场景、人物特征、文字、行为、背景和任何其他重要细节。"
                        }
                    ],
                },
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "请详细分析这张图片的内容:"
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/{ext[1:]};base64,{encoded_image}"
                            }
                        }
                    ],
                }
            ]
            
            # 调用模型进行图像分析
            response = self.model(messages)
            
            # 返回分析结果
            return response.content if hasattr(response, 'content') else str(response)
                
        except Exception as e:
            return f"分析图像时发生错误: {str(e)}"





